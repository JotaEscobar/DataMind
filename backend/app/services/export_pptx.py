"""
export_pptx.py — Exportación PPTX con gráficos reales y branding DataMind

Genera un deck ejecutivo con:
  - Slide de portada con branding DataMind
  - Slide de métricas clave (KPIs)
  - Slides de gráficos (bar, line, pie generados con matplotlib)
  - Slide de hallazgos e insights
  - Slide de próximos pasos

Uso:
    from app.tools.export_pptx import build_pptx_report
    path = build_pptx_report(file_path, session_id, title, insights)
"""
from __future__ import annotations

import io
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Inches, Pt
from lxml import etree

# ── DataMind palette ──────────────────────────────────────────────────────────
_C = RGBColor

DM_DARK    = _C(0x08, 0x0d, 0x14)
DM_SURFACE = _C(0x0e, 0x17, 0x24)
DM_BORDER  = _C(0x1a, 0x2d, 0x42)
DM_ACCENT  = _C(0x00, 0xd4, 0xff)
DM_GREEN   = _C(0x00, 0xe5, 0xa0)
DM_AMBER   = _C(0xff, 0xb3, 0x00)
DM_RED     = _C(0xff, 0x45, 0x60)
DM_PURPLE  = _C(0xa8, 0x55, 0xf7)
DM_WHITE   = _C(0xf0, 0xf4, 0xf8)
DM_MID     = _C(0x7a, 0x9a, 0xb8)
DM_DIM     = _C(0x4a, 0x68, 0x85)

CHART_COLORS = [
    "#00d4ff", "#00e5a0", "#ffb300", "#a855f7", "#ff4560",
    "#0066ff", "#ff6b6b", "#ffd93d", "#6bcb77", "#4d96ff",
]

STORAGE_DIR = Path(__file__).resolve().parents[3] / "data_storage" / "exports"

# Slide size 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Matplotlib chart → BytesIO ────────────────────────────────────────────────

def _setup_chart_style():
    plt.rcParams.update({
        "figure.facecolor": "#111d2e",
        "axes.facecolor":   "#0e1724",
        "axes.edgecolor":   "#1a2d42",
        "axes.labelcolor":  "#7a9ab8",
        "xtick.color":      "#4a6885",
        "ytick.color":      "#4a6885",
        "text.color":       "#f0f4f8",
        "grid.color":       "#1a2d42",
        "font.family":      "sans-serif",
        "font.size":        10,
    })


def _chart_to_stream(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _bar_chart_stream(
    labels: List[str], values: List[float], title: str,
    ylabel: str = "", horizontal: bool = False, color: str = "#00d4ff",
) -> io.BytesIO:
    _setup_chart_style()
    fig, ax = plt.subplots(figsize=(11, 5))
    fig.patch.set_facecolor("#111d2e")
    labels = [str(l)[:22] for l in labels]

    if horizontal:
        bars = ax.barh(labels, values, color=color, alpha=0.85, height=0.55)
        for bar, val in zip(bars, values):
            ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                    f"{val:,.0f}", va="center", fontsize=8.5, color="#f0f4f8")
    else:
        bars = ax.bar(range(len(labels)), values, color=color, alpha=0.85, width=0.6)
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=8.5)
        if ylabel: ax.set_ylabel(ylabel, color="#7a9ab8")
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() * 1.01,
                    f"{val:,.0f}", ha="center", fontsize=7.5, color="#f0f4f8")

    ax.set_title(title, color="#f0f4f8", fontsize=12, pad=12)
    ax.grid(axis="x" if horizontal else "y", alpha=0.25)
    ax.spines[:].set_visible(False)
    fig.tight_layout()
    return _chart_to_stream(fig)


def _line_chart_stream(
    x_vals: List, y_vals: List[float], title: str,
    x_label: str = "", y_label: str = "", color: str = "#00e5a0",
) -> io.BytesIO:
    _setup_chart_style()
    fig, ax = plt.subplots(figsize=(11, 5))
    fig.patch.set_facecolor("#111d2e")
    ax.plot(range(len(x_vals)), y_vals, color=color, linewidth=2.5,
            marker="o", markersize=4)
    ax.fill_between(range(len(x_vals)), y_vals, alpha=0.12, color=color)
    step = max(1, len(x_vals) // 8)
    ax.set_xticks(range(0, len(x_vals), step))
    ax.set_xticklabels([str(x_vals[i])[:10] for i in range(0, len(x_vals), step)],
                       rotation=25, ha="right", fontsize=8.5)
    ax.set_title(title, color="#f0f4f8", fontsize=12, pad=12)
    if x_label: ax.set_xlabel(x_label, color="#7a9ab8")
    if y_label: ax.set_ylabel(y_label, color="#7a9ab8")
    ax.grid(axis="y", alpha=0.25)
    ax.spines[:].set_visible(False)
    fig.tight_layout()
    return _chart_to_stream(fig)


def _pie_chart_stream(labels: List[str], values: List[float], title: str) -> io.BytesIO:
    _setup_chart_style()
    if len(labels) > 7:
        labels, values = labels[:6] + ["Otros"], values[:6] + [sum(values[6:])]
    fig, ax = plt.subplots(figsize=(8, 5))
    fig.patch.set_facecolor("#111d2e")
    wedges, _, autotexts = ax.pie(
        values, autopct="%1.1f%%",
        colors=CHART_COLORS[:len(values)], startangle=90,
        wedgeprops={"edgecolor": "#0e1724", "linewidth": 1.5},
        pctdistance=0.82,
    )
    for at in autotexts:
        at.set_color("#f0f4f8")
        at.set_fontsize(9)
    ax.legend(wedges, [str(l)[:18] for l in labels],
              loc="lower center", bbox_to_anchor=(0.5, -0.18),
              ncol=3, fontsize=8, frameon=False, labelcolor="#7a9ab8")
    ax.set_title(title, color="#f0f4f8", fontsize=12, pad=10)
    fig.tight_layout()
    return _chart_to_stream(fig)


# ── PPTX helpers ───────────────────────────────────────────────────────────────

def _set_bg(slide, rgb: RGBColor):
    """Fondo sólido para un slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_text_box(
    slide, text: str, left: float, top: float,
    width: float, height: float,
    font_size: int = 14, bold: bool = False,
    color: RGBColor = DM_WHITE, align=PP_ALIGN.LEFT,
    italic: bool = False,
) -> Any:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _add_rect(slide, left: float, top: float, width: float, height: float,
              fill: RGBColor, line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def _add_chart_image(slide, stream: io.BytesIO,
                     left: float, top: float, width: float, height: float):
    slide.shapes.add_picture(stream, Inches(left), Inches(top),
                              Inches(width), Inches(height))


# ── Slides ─────────────────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, title: str, file_name: str, df: Optional[pd.DataFrame]):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DM_DARK)

    # Franja izquierda cyan
    _add_rect(slide, 0, 0, 0.08, 7.5, DM_ACCENT)

    # Bloque superior oscuro
    _add_rect(slide, 0.08, 0, 13.25, 3.5, _C(0x0e, 0x17, 0x24))

    # Etiqueta producto
    _add_text_box(slide, "DATAMIND 2.0", 0.5, 0.3, 8, 0.5,
                  font_size=10, color=DM_ACCENT, bold=True)

    # Título principal
    _add_text_box(slide, title or "Reporte de Análisis de Datos",
                  0.5, 0.9, 12, 1.8,
                  font_size=32, bold=True, color=DM_WHITE)

    # Fecha y archivo
    date_str = datetime.now().strftime("%d de %B de %Y")
    meta = f"{date_str}  ·  {file_name}"
    _add_text_box(slide, meta, 0.5, 2.8, 10, 0.5,
                  font_size=11, color=DM_MID)

    # Línea separadora
    _add_rect(slide, 0.5, 3.55, 12.8, 0.02, DM_ACCENT)

    # KPIs en la parte baja
    if df is not None:
        num_df = df.select_dtypes(include="number")
        total_missing = int(df.isnull().sum().sum())
        quality = round(100 - total_missing / max(len(df) * len(df.columns), 1) * 100, 1)

        kpis = [
            (f"{len(df):,}", "REGISTROS"),
            (f"{len(df.columns)}", "VARIABLES"),
            (f"{quality}%", "CALIDAD"),
            (f"{len(num_df.columns)}", "COL. NUMÉRICAS"),
        ]
        box_w = 2.8
        for i, (val, label) in enumerate(kpis):
            x = 0.5 + i * (box_w + 0.2)
            _add_rect(slide, x, 4.0, box_w, 2.2, _C(0x11, 0x1d, 0x2e),
                      line=DM_BORDER)
            _add_text_box(slide, val, x, 4.25, box_w, 1.0,
                          font_size=28, bold=True, color=DM_ACCENT,
                          align=PP_ALIGN.CENTER)
            _add_text_box(slide, label, x, 5.3, box_w, 0.5,
                          font_size=9, color=DM_MID, align=PP_ALIGN.CENTER)

    # Branding footer
    _add_text_box(slide, "Reporte Confidencial — Generado por DataMind AI",
                  0.5, 7.1, 12, 0.3, font_size=8, color=DM_DIM)


def _slide_insights(prs: Presentation, insights: str, warnings: Optional[List[str]]):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DM_DARK)
    _add_rect(slide, 0, 0, 0.06, 7.5, DM_ACCENT)

    _add_text_box(slide, "Hallazgos Principales", 0.4, 0.2, 12, 0.7,
                  font_size=22, bold=True, color=DM_ACCENT)
    _add_rect(slide, 0.4, 1.0, 12.8, 0.02, _C(0x1a, 0x2d, 0x42))

    # Limpiar markdown básico
    clean = (insights or "").replace("**", "").strip()
    paragraphs = [p.strip() for p in clean.split("\n\n") if p.strip()][:4]

    y = 1.1
    for para in paragraphs:
        # Truncar párrafos muy largos
        if len(para) > 300:
            para = para[:297] + "…"
        box_h = 0.9 + para.count("\n") * 0.2
        _add_rect(slide, 0.4, y, 12.4, box_h, _C(0x0e, 0x17, 0x24),
                  line=_C(0x1a, 0x2d, 0x42))
        _add_text_box(slide, para, 0.6, y + 0.1, 12.0, box_h - 0.15,
                      font_size=11, color=_C(0xe2, 0xe8, 0xf0))
        y += box_h + 0.15
        if y > 6.5:
            break

    if warnings:
        _add_text_box(slide, "⚠  " + "  ·  ".join(warnings[:3]),
                      0.4, 6.9, 12, 0.45,
                      font_size=9, color=DM_AMBER, italic=True)


def _slide_chart(prs: Presentation, chart_stream: io.BytesIO,
                 slide_title: str, insight_line: str = ""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DM_DARK)
    _add_rect(slide, 0, 0, 0.06, 7.5, DM_GREEN)

    _add_text_box(slide, slide_title, 0.4, 0.15, 12.5, 0.65,
                  font_size=18, bold=True, color=DM_WHITE)
    _add_rect(slide, 0.4, 0.85, 12.8, 0.02, _C(0x1a, 0x2d, 0x42))

    _add_chart_image(slide, chart_stream, 0.4, 1.0, 12.5, 5.8)

    if insight_line:
        _add_rect(slide, 0.4, 6.85, 12.4, 0.52, _C(0x0e, 0x17, 0x24),
                  line=_C(0x00, 0xd4, 0xff))
        _add_text_box(slide, f"💡 {insight_line[:120]}", 0.6, 6.9, 12.0, 0.42,
                      font_size=9.5, color=DM_ACCENT)


def _slide_next_steps(prs: Presentation, questions: List[str]):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, DM_DARK)
    _add_rect(slide, 0, 0, 0.06, 7.5, DM_PURPLE)

    _add_text_box(slide, "Análisis Recomendados", 0.4, 0.2, 12, 0.7,
                  font_size=22, bold=True, color=DM_PURPLE)
    _add_rect(slide, 0.4, 1.0, 12.8, 0.02, _C(0x1a, 0x2d, 0x42))

    colors_cycle = [DM_ACCENT, DM_GREEN, DM_AMBER, DM_PURPLE]
    for i, q in enumerate(questions[:6]):
        y = 1.2 + i * 0.95
        accent = colors_cycle[i % len(colors_cycle)]
        _add_rect(slide, 0.4, y, 0.05, 0.7, accent)
        _add_rect(slide, 0.5, y, 12.3, 0.7, _C(0x0e, 0x17, 0x24),
                  line=_C(0x1a, 0x2d, 0x42))
        _add_text_box(slide, f"{i+1}.  {q}", 0.65, y + 0.08, 12.0, 0.55,
                      font_size=12, color=DM_WHITE)

    _add_text_box(slide, "Continúa el análisis en DataMind AI",
                  0.4, 7.1, 12, 0.3, font_size=8, color=DM_DIM)


# ── Main builder ───────────────────────────────────────────────────────────────

def build_pptx_report(
    file_path: str,
    session_id: str,
    title: str,
    insights: str,
    warnings: Optional[List[str]] = None,
    suggested_questions: Optional[List[str]] = None,
) -> str:
    """
    Construye el PPTX completo y devuelve la ruta del archivo generado.
    """
    STORAGE_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = str(STORAGE_DIR / f"datamind_{session_id[:8]}_{timestamp}.pptx")

    df = _load_df(file_path)
    file_name = Path(file_path).name

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Portada
    _slide_cover(prs, title, file_name, df)

    # 2. Hallazgos
    _slide_insights(prs, insights, warnings)

    # 3. Gráficos (hasta 3)
    if df is not None:
        num_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        # Gráfico 1: top categorías
        if cat_cols and num_cols:
            cat_col, num_col = cat_cols[0], num_cols[0]
            grouped = (df.groupby(cat_col)[num_col].sum()
                       .sort_values(ascending=False).head(10))
            if len(grouped) > 1:
                stream = _bar_chart_stream(
                    grouped.index.tolist(), grouped.values.tolist(),
                    f"Top 10 — {num_col} por {cat_col}",
                    ylabel=num_col, horizontal=len(grouped) > 5,
                    color="#00d4ff",
                )
                top_val = grouped.iloc[0]
                top_label = grouped.index[0]
                insight_line = (
                    f"{top_label} lidera con {top_val:,.0f} "
                    f"({top_val/grouped.sum()*100:.1f}% del total)"
                )
                _slide_chart(prs, stream,
                             f"Distribución de {num_col}",
                             insight_line)

        # Gráfico 2: serie temporal
        date_col = _detect_date_col(df)
        if date_col and num_cols:
            try:
                ts_df = df.copy()
                ts_df[date_col] = pd.to_datetime(ts_df[date_col], errors="coerce")
                ts_df = ts_df.dropna(subset=[date_col]).set_index(date_col).sort_index()
                num_col = num_cols[0]
                monthly = ts_df[num_col].resample("ME").sum()
                if len(monthly) >= 3:
                    stream = _line_chart_stream(
                        [d.strftime("%b %Y") for d in monthly.index],
                        monthly.values.tolist(),
                        f"Tendencia mensual — {num_col}",
                        y_label=num_col, color="#00e5a0",
                    )
                    pct = ((monthly.iloc[-1] - monthly.iloc[0]) /
                           abs(monthly.iloc[0]) * 100) if monthly.iloc[0] != 0 else 0
                    trend = "creció" if pct > 0 else "cayó"
                    insight_line = (
                        f"{num_col} {trend} un {abs(pct):.1f}% "
                        f"de {monthly.index[0].strftime('%b %Y')} "
                        f"a {monthly.index[-1].strftime('%b %Y')}"
                    )
                    _slide_chart(prs, stream,
                                 f"Evolución temporal — {num_col}",
                                 insight_line)
            except Exception:
                pass

        # Gráfico 3: donut categorías
        if cat_cols and num_cols:
            try:
                cat_col, num_col = cat_cols[0], num_cols[0]
                grouped = (df.groupby(cat_col)[num_col].sum()
                           .sort_values(ascending=False).head(7))
                if 2 <= len(grouped) <= 10:
                    stream = _pie_chart_stream(
                        grouped.index.tolist(), grouped.values.tolist(),
                        f"Composición — {num_col} por {cat_col}",
                    )
                    top_two_pct = grouped.iloc[:2].sum() / grouped.sum() * 100
                    insight_line = (
                        f"Los 2 primeros segmentos concentran "
                        f"el {top_two_pct:.1f}% del total"
                    )
                    _slide_chart(prs, stream,
                                 f"Composición por {cat_col}",
                                 insight_line)
            except Exception:
                pass

    # 4. Próximos pasos
    if suggested_questions:
        _slide_next_steps(prs, suggested_questions)

    prs.save(out_path)
    return out_path


def _load_df(file_path: str) -> Optional[pd.DataFrame]:
    try:
        ext = Path(file_path).suffix.lower()
        if ext == ".csv":
            return pd.read_csv(file_path)
        elif ext in {".xlsx", ".xls"}:
            return pd.read_excel(file_path)
    except Exception:
        pass
    return None


def _detect_date_col(df: pd.DataFrame) -> Optional[str]:
    for col in df.select_dtypes(include=["datetime", "datetimetz"]).columns:
        return col
    for col in df.select_dtypes(include="object").columns:
        try:
            parsed = pd.to_datetime(df[col], errors="coerce")
            if parsed.notnull().mean() > 0.7:
                df[col] = parsed
                return col
        except Exception:
            pass
    return None